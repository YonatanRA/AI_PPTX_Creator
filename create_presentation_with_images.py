#!/usr/bin/env python3
"""
AI PPTX Creator with Gemini Image Generation (Nano Banana)

This script generates PowerPoint presentations with AI-generated images using Google's
Gemini API (Imagen model, internally codenamed "Nano Banana"). It creates visually
compelling slides based on PDF document context.
"""

import os
import argparse
import logging
from pathlib import Path
from typing import List, Optional
from io import BytesIO

from dotenv import load_dotenv
from google import genai
from google.genai import types
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from langchain_openai.chat_models import ChatOpenAI
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_community.vectorstores import Chroma
from langchain.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class PresentationConfig:
    """Configuration for presentation generation with image support."""
    
    def __init__(
        self,
        pdf_directory: str = "pdfs",
        output_directory: str = "pptx",
        chroma_directory: str = "chroma_db",
        images_directory: str = "imgs/generated",
        openai_model: str = "gpt-4-turbo",
        imagen_model: str = "imagen-3.0-generate-002",
        num_slides: int = 5,
        bullet_points_per_slide: int = 3,
        bullet_point_words: int = 30,
        retriever_k: int = 3,
        retriever_lambda: float = 0.25,
        image_aspect_ratio: str = "16:9"
    ):
        """Initialize presentation configuration.
        
        Args:
            pdf_directory: Directory containing PDF files
            output_directory: Directory for output presentations
            chroma_directory: Directory for ChromaDB storage
            images_directory: Directory for generated images
            openai_model: OpenAI model name for text generation
            imagen_model: Imagen model for image generation
            num_slides: Number of content slides to generate
            bullet_points_per_slide: Number of bullet points per slide
            bullet_point_words: Maximum words per bullet point
            retriever_k: Number of documents to retrieve
            retriever_lambda: Lambda multiplier for MMR algorithm
            image_aspect_ratio: Aspect ratio for generated images
        """
        self.pdf_directory = Path(pdf_directory)
        self.output_directory = Path(output_directory)
        self.chroma_directory = Path(chroma_directory)
        self.images_directory = Path(images_directory)
        self.openai_model = openai_model
        self.imagen_model = imagen_model
        self.num_slides = num_slides
        self.bullet_points_per_slide = bullet_points_per_slide
        self.bullet_point_words = bullet_point_words
        self.retriever_k = retriever_k
        self.retriever_lambda = retriever_lambda
        self.image_aspect_ratio = image_aspect_ratio
        
        # Ensure directories exist
        self.output_directory.mkdir(parents=True, exist_ok=True)
        self.chroma_directory.mkdir(parents=True, exist_ok=True)
        self.images_directory.mkdir(parents=True, exist_ok=True)


class AIImagePresentationCreator:
    """Main class for creating AI-powered presentations with generated images."""
    
    def __init__(self, config: PresentationConfig):
        """Initialize the presentation creator.
        
        Args:
            config: Configuration object for presentation generation
        """
        self.config = config
        self._load_api_keys()
        
        # Initialize OpenAI components
        self.chat_model = ChatOpenAI(model=config.openai_model)
        self.embeddings = OpenAIEmbeddings()
        self.parser = StrOutputParser()
        
        # Initialize Google Gemini client for image generation
        self.gemini_client = genai.Client(api_key=self.google_api_key)
        
        self.vector_db = None
        self.retriever = None
        
    def _load_api_keys(self):
        """Load API keys from environment."""
        load_dotenv()
        
        self.openai_api_key = os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError(
                "OPENAI_API_KEY not found. Please create a .env file with your API key."
            )
        
        # Google API key can be GOOGLE_API_KEY or GEMINI_API_KEY
        self.google_api_key = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
        if not self.google_api_key:
            raise ValueError(
                "GOOGLE_API_KEY or GEMINI_API_KEY not found. Please add it to your .env file."
            )
        
        logger.info("API keys loaded successfully")
    
    def load_documents(self) -> List:
        """Load PDF documents from the configured directory.
        
        Returns:
            List of loaded document pages
        """
        logger.info(f"Loading PDF documents from {self.config.pdf_directory}")
        
        if not self.config.pdf_directory.exists():
            raise FileNotFoundError(
                f"PDF directory not found: {self.config.pdf_directory}"
            )
        
        loader = PyPDFDirectoryLoader(str(self.config.pdf_directory))
        pages = loader.load()
        
        logger.info(f"Loaded {len(pages)} pages from PDF documents")
        return pages
    
    def create_vector_database(self, documents: List, force_recreate: bool = False):
        """Create or load vector database from documents.
        
        Args:
            documents: List of documents to embed
            force_recreate: If True, recreate the database even if it exists
        """
        if force_recreate or not self.config.chroma_directory.exists():
            logger.info("Creating new vector database")
            self.vector_db = Chroma.from_documents(
                documents,
                self.embeddings,
                persist_directory=str(self.config.chroma_directory)
            )
        else:
            logger.info("Loading existing vector database")
            self.vector_db = Chroma(
                persist_directory=str(self.config.chroma_directory),
                embedding_function=self.embeddings
            )
        
        # Create retriever with MMR for diverse results
        self.retriever = self.vector_db.as_retriever(
            search_type="mmr",
            search_kwargs={
                "k": self.config.retriever_k,
                "lambda_mult": self.config.retriever_lambda
            }
        )
        logger.info("Vector database and retriever initialized")
    
    def generate_slide_content(self, query: str) -> dict:
        """Generate structured slide content from a query using RAG.
        
        Args:
            query: Query to generate slide content for
            
        Returns:
            Dictionary containing title and slides with bullet points
        """
        logger.info(f"Generating slide content for query: {query}")
        
        template = f"""
            Given the context below and the question, 
            please generate a presentation structure with:
            - A main title for the presentation
            - {self.config.num_slides} slides, each with:
              * A slide title
              * {self.config.bullet_points_per_slide} bullet points
              * Each bullet point should be {self.config.bullet_point_words} words maximum
            
            Format your response as follows:
            PRESENTATION TITLE: [Main title here]
            
            SLIDE 1: [Slide title]
            - [Bullet point 1]
            - [Bullet point 2]
            - [Bullet point 3]
            
            SLIDE 2: [Slide title]
            - [Bullet point 1]
            - [Bullet point 2]
            - [Bullet point 3]
            
            And so on...

            Context: {{context}}
            Question: {{question}}
            """
        
        prompt = ChatPromptTemplate.from_template(template)
        
        # Create RAG chain
        chain = (
            {"context": self.retriever, "question": RunnablePassthrough()}
            | prompt
            | self.chat_model
            | self.parser
        )
        
        response = chain.invoke(query)
        logger.info("Slide content generated successfully")
        
        return self._parse_slide_content(response)
    
    def _parse_slide_content(self, content: str) -> dict:
        """Parse the LLM response into structured slide data.
        
        Args:
            content: Raw text response from LLM
            
        Returns:
            Structured dictionary with presentation data
        """
        lines = content.strip().split('\n')
        presentation = {
            'title': '',
            'slides': []
        }
        
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('PRESENTATION TITLE:'):
                presentation['title'] = line.replace('PRESENTATION TITLE:', '').strip()
            elif line.startswith('SLIDE '):
                if current_slide:
                    presentation['slides'].append(current_slide)
                # Extract slide title
                title = line.split(':', 1)[1].strip() if ':' in line else line
                current_slide = {'title': title, 'bullets': []}
            elif line.startswith('-') or line.startswith('•'):
                if current_slide:
                    bullet = line.lstrip('-•').strip()
                    current_slide['bullets'].append(bullet)
        
        # Add last slide
        if current_slide:
            presentation['slides'].append(current_slide)
        
        return presentation
    
    def generate_image_for_slide(self, slide_title: str, slide_content: List[str]) -> Optional[Path]:
        """Generate an AI image for a slide using Gemini Imagen.
        
        Args:
            slide_title: Title of the slide
            slide_content: List of bullet points for context
            
        Returns:
            Path to the generated image file, or None if generation fails
        """
        logger.info(f"Generating image for slide: {slide_title}")
        
        # Create a descriptive prompt for image generation
        content_summary = ' '.join(slide_content[:2])  # Use first 2 bullet points
        prompt = f"Professional presentation illustration for: {slide_title}. Context: {content_summary}. Modern, clean, business style."
        
        try:
            response = self.gemini_client.models.generate_images(
                model=self.config.imagen_model,
                prompt=prompt,
                config=types.GenerateImagesConfig(
                    number_of_images=1,
                    aspect_ratio=self.config.image_aspect_ratio,
                    safety_filter_level="BLOCK_ONLY_HIGH",
                )
            )
            
            if response.generated_images:
                # Save the image
                image_bytes = response.generated_images[0].image.image_bytes
                image = Image.open(BytesIO(image_bytes))
                
                # Create filename from slide title
                filename = "".join(c for c in slide_title if c.isalnum() or c in (' ', '-', '_')).strip()
                filename = filename[:50]  # Limit length
                image_path = self.config.images_directory / f"{filename}.png"
                
                image.save(image_path)
                logger.info(f"Image saved to {image_path}")
                return image_path
            else:
                logger.warning(f"No images generated for: {slide_title}")
                return None
                
        except Exception as e:
            logger.error(f"Error generating image for '{slide_title}': {e}")
            return None
    
    def create_presentation_file(
        self,
        presentation_data: dict,
        output_filename: str,
        generate_images: bool = True
    ):
        """Create a PowerPoint file with the generated content and images.
        
        Args:
            presentation_data: Structured presentation data
            output_filename: Name of the output file
            generate_images: Whether to generate AI images for slides
        """
        logger.info("Creating PowerPoint presentation")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = presentation_data.get('title', 'AI-Generated Presentation')
        
        # Content slides
        for slide_data in presentation_data['slides']:
            # Use blank layout for custom positioning
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
            )
            title_frame = title_shape.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            
            # Generate and add image if enabled
            image_height = 0
            if generate_images:
                image_path = self.generate_image_for_slide(
                    slide_data['title'],
                    slide_data['bullets']
                )
                if image_path and image_path.exists():
                    # Add image to the right side
                    try:
                        pic = slide.shapes.add_picture(
                            str(image_path),
                            Inches(5.5), Inches(1.5),
                            width=Inches(4), height=Inches(3)
                        )
                        image_height = 3
                    except Exception as e:
                        logger.warning(f"Could not add image to slide: {e}")
            
            # Add bullet points on the left
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(slide_data['bullets']):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.space_after = Pt(12)
        
        # Save presentation
        output_path = self.config.output_directory / output_filename
        prs.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")
    
    def create_presentation(
        self,
        query: str,
        output_filename: str,
        force_recreate_db: bool = False,
        generate_images: bool = True
    ):
        """Complete workflow to create a presentation with AI-generated images.
        
        Args:
            query: Query to generate presentation content for
            output_filename: Name of the output PowerPoint file
            force_recreate_db: If True, recreate the vector database
            generate_images: Whether to generate AI images for slides
        """
        # Load documents and create vector database
        documents = self.load_documents()
        self.create_vector_database(documents, force_recreate=force_recreate_db)
        
        # Generate slide content
        presentation_data = self.generate_slide_content(query)
        
        # Create PowerPoint file with images
        self.create_presentation_file(
            presentation_data,
            output_filename,
            generate_images=generate_images
        )
        
        logger.info("✅ Presentation creation complete!")


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentations with AI-generated images using Gemini"
    )
    parser.add_argument(
        "--query",
        type=str,
        default="What are the main topics and key points of this document?",
        help="Query to generate presentation content for"
    )
    parser.add_argument(
        "--output",
        type=str,
        default="AI_Presentation_with_Images.pptx",
        help="Name of the output PowerPoint file"
    )
    parser.add_argument(
        "--pdf-dir",
        type=str,
        default="pdfs",
        help="Directory containing PDF files"
    )
    parser.add_argument(
        "--output-dir",
        type=str,
        default="pptx",
        help="Directory for output presentations"
    )
    parser.add_argument(
        "--num-slides",
        type=int,
        default=5,
        help="Number of content slides to generate"
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Disable AI image generation (text-only slides)"
    )
    parser.add_argument(
        "--force-recreate-db",
        action="store_true",
        help="Force recreation of the vector database"
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose logging"
    )
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Create configuration
    config = PresentationConfig(
        pdf_directory=args.pdf_dir,
        output_directory=args.output_dir,
        num_slides=args.num_slides
    )
    
    # Create presentation
    creator = AIImagePresentationCreator(config)
    creator.create_presentation(
        query=args.query,
        output_filename=args.output,
        force_recreate_db=args.force_recreate_db,
        generate_images=not args.no_images
    )


if __name__ == "__main__":
    main()
