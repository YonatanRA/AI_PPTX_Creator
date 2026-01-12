import os
import sys
import logging
from io import BytesIO
from typing import List, Dict

# Third-party imports
from dotenv import load_dotenv
from langchain_openai.chat_models import ChatOpenAI
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from google import genai
from google.genai import types
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Load environment variables
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY not found in environment variables.")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found. Please add it to your .env file.")

class NanobananaPresentationCreator:
    def __init__(self):
        self.openai_model = "gpt-4-turbo"
        # "Nano Banana" refers to the Imagen 3 model in this context
        self.nanobanana_model = "imagen-3.0-generate-002" 
        
        self.embeddings = OpenAIEmbeddings()
        self.chat_model = ChatOpenAI(model=self.openai_model)
        self.gemini_client = genai.Client(api_key=GOOGLE_API_KEY)
        self.parser = StrOutputParser()
        
        self.vector_db = None
        self.retriever = None

    def load_and_process_documents(self, pdf_dir: str = "pdfs", chroma_dir: str = "chroma_db"):
        """Loads PDFs and creates/loads a vector database."""
        if not os.path.exists(pdf_dir):
            raise FileNotFoundError(f"Directory {pdf_dir} does not exist.")

        logging.info(f"Loading PDFs from {pdf_dir}...")
        loader = PyPDFDirectoryLoader(pdf_dir)
        pages = loader.load()
        logging.info(f"Loaded {len(pages)} pages.")

        logging.info("Initializing Vector Database...")
        # Using existing persist directory if available, else create new
        self.vector_db = Chroma.from_documents(
            pages, 
            self.embeddings, 
            persist_directory=chroma_dir
        )
        
        self.retriever = self.vector_db.as_retriever(
            search_type="mmr", 
            search_kwargs={"k": 2, "lambda_mult": 0.25}
        )
        logging.info("Vector Database ready.")

    def generate_presentation_plan(self, query: str) -> Dict:
        """Generates a structured plan for the presentation (Titles & Text)."""
        logging.info("Generating presentation plan...")
        
        template = """
        Given the context below and the user's request, create a plan for a 3-slide presentation. 
        
        Output strictly in this format (no markdown code blocks):
        TITLE: <Presentation Title>
        SLIDE 1: <Slide Title> | <Bullet 1> | <Bullet 2> | <Bullet 3>
        SLIDE 2: <Slide Title> | <Bullet 1> | <Bullet 2> | <Bullet 3>
        SLIDE 3: <Slide Title> | <Bullet 1> | <Bullet 2> | <Bullet 3>

        Context: {context}
        Request: {question}
        """
        
        prompt = ChatPromptTemplate.from_template(template)
        chain = (
            {"context": self.retriever, "question": RunnablePassthrough()}
            | prompt
            | self.chat_model
            | self.parser
        )
        
        response = chain.invoke(query)
        logging.info("Plan generated.")
        return self._parse_plan(response)

    def _parse_plan(self, response_text: str) -> Dict:
        """Parses the LLM output into a dictionary."""
        lines = response_text.strip().split('\n')
        plan = {'title': 'Presentation', 'slides': []}
        
        for line in lines:
            line = line.strip()
            if line.startswith("TITLE:"):
                plan['title'] = line.replace("TITLE:", "").strip()
            elif line.startswith("SLIDE"):
                parts = line.split(":", 1)[1].split("|")
                slide = {
                    'title': parts[0].strip(),
                    'bullets': [p.strip() for p in parts[1:] if p.strip()]
                }
                plan['slides'].append(slide)
        return plan

    def generate_nanobanana_image(self, prompt: str, output_path: str):
        """Generates an image using Google Nanobanana (Imagen 3)."""
        logging.info(f"Generating image for: {prompt[:50]}...")
        try:
            response = self.gemini_client.models.generate_images(
                model=self.nanobanana_model,
                prompt=f"Professional presentation illustration, modern, high quality: {prompt}",
                config=types.GenerateImagesConfig(
                    number_of_images=1,
                    aspect_ratio="16:9", # Standard PPTX aspect ratio
                    safety_filter_level="BLOCK_ONLY_HIGH",
                )
            )
            
            if response.generated_images:
                image_bytes = response.generated_images[0].image.image_bytes
                image = Image.open(BytesIO(image_bytes))
                image.save(output_path)
                return True
        except Exception as e:
            logging.error(f"Failed to generate image: {e}")
            return False

    def create_presentation(self, query: str, output_file: str):
        """Orchestrates the creation of the PPTX file."""
        self.load_and_process_documents()
        plan = self.generate_presentation_plan(query)
        
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = plan['title']
        subtitle.text = "Generated by AI (Nanobanana Edition)"
        
        # Content Slides
        img_dir = "imgs/generated_nanobanana"
        os.makedirs(img_dir, exist_ok=True)
        
        for i, slide_data in enumerate(plan['slides']):
            logging.info(f"Processing Slide {i+1}: {slide_data['title']}")
            
            # Use a blank layout to freely place image and text
            blank_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_layout)
            
            # 1. Add Title
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = slide_data['title']
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True
            
            # 2. Generate and Add Image (Nanobanana)
            # Create a prompt based on title and bullets
            image_prompt = f"{slide_data['title']}. " + " ".join(slide_data['bullets'])
            img_filename = f"slide_{i+1}.png"
            img_path = os.path.join(img_dir, img_filename)
            
            if self.generate_nanobanana_image(image_prompt, img_path):
                # Place image on the right
                img_left = Inches(5.5)
                img_top = Inches(2)
                img_width = Inches(4)
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
            
            # 3. Add Text (Bullets) on the left
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(4.5)
            height = Inches(5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for bullet in slide_data['bullets']:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.level = 0
                
        output_path = os.path.join("pptx", output_file)
        os.makedirs("pptx", exist_ok=True)
        prs.save(output_path)
        logging.info(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    creator = NanobananaPresentationCreator()
    creator.create_presentation(
        query="What are the key security threats in the Red Sea?",
        output_file="Red_Sea_Nanobanana.pptx"
    )
